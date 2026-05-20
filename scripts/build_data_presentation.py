#!/usr/bin/env python3
"""Build a concise 10-slide technical PPTX summarising the 6 notebooks.

Audience: academic supervisor. Tone: technical, bilingual (FR + ML EN terms).
Each slide maps to one notebook (or a synthesis slide).

Output: presentation/movie_success_technical_deck.pptx
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
SRC = ROOT / "src"
if str(SRC) not in sys.path:
    sys.path.insert(0, str(SRC))

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import config

PRESENTATION_DIR = ROOT / "presentation"
OUTPUT_PPTX = PRESENTATION_DIR / "movie_success_technical_deck.pptx"

# Design system (aligned with the Streamlit app)
PRIMARY = RGBColor(0x1E, 0x3A, 0x8A)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
ACCENT_SOFT = RGBColor(0xDB, 0xEA, 0xFE)
DARK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x47, 0x55, 0x69)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS = RGBColor(0x15, 0x80, 0x3D)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Header layout (inches)
MARGIN_X = 0.55
CONTENT_WIDTH = 13.333 - (2 * MARGIN_X)
TITLE_TOP = 0.35
SUBTITLE_TOP = 0.92
BODY_TOP = 1.45


class TechDeck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.slide_num = 0
        self.plots_used: list[str] = []
        self.plots_missing: list[str] = []
        self.kpis = self._load_json(config.APP_KPIS_FILE)
        self.regime = self._load_csv(config.REGIME_COMPARISON_FILE)
        self.metrics = self._load_csv_first_row(config.MODEL_METRICS_FILE)

    @staticmethod
    def _load_json(path: Path) -> dict:
        if not path.exists():
            return {}
        with path.open(encoding="utf-8") as fh:
            return json.load(fh)

    @staticmethod
    def _load_csv(path: Path) -> pd.DataFrame:
        return pd.read_csv(path) if path.exists() else pd.DataFrame()

    @staticmethod
    def _load_csv_first_row(path: Path) -> dict:
        if not path.exists():
            return {}
        df = pd.read_csv(path)
        return df.iloc[0].to_dict() if len(df) else {}

    # ----- shape helpers -----
    def _plot(self, name: str, subdir: str) -> Path:
        return config.PLOTS_DIR / subdir / name

    def _pct(self, key: str, fallback: float = 0.0) -> str:
        val = self.kpis.get(key, self.metrics.get(key, fallback))
        try:
            return f"{float(val):.1%}"
        except (TypeError, ValueError):
            return "—"

    def _num(self, key: str, fallback: int = 0) -> str:
        val = self.kpis.get(key, fallback)
        try:
            return f"{int(val):,}"
        except (TypeError, ValueError):
            return str(fallback)

    def _new_slide(self):
        self.slide_num += 1
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        # background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.fill.background()
        # footer slide number
        foot = slide.shapes.add_textbox(Inches(12.3), Inches(7.1), Inches(0.8), Inches(0.3))
        foot.line.fill.background()
        p = foot.text_frame.paragraphs[0]
        p.text = f"{self.slide_num} / 10"
        p.font.size = Pt(10)
        p.font.color.rgb = MUTED
        p.alignment = PP_ALIGN.RIGHT
        return slide

    def _accent_bar(self, slide) -> None:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.0), SLIDE_W, Inches(0.18))
        bar.fill.solid()
        bar.fill.fore_color.rgb = PRIMARY
        bar.line.fill.background()

    def _title(self, slide, title: str, subtitle: str | None = None) -> None:
        """Title + subtitle across the full slide width."""
        tb = slide.shapes.add_textbox(
            Inches(MARGIN_X), Inches(TITLE_TOP), Inches(CONTENT_WIDTH), Inches(0.85)
        )
        tb.line.fill.background()
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        if subtitle:
            tb2 = slide.shapes.add_textbox(
                Inches(MARGIN_X), Inches(SUBTITLE_TOP), Inches(CONTENT_WIDTH), Inches(0.42)
            )
            tb2.line.fill.background()
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(12)
            p2.font.color.rgb = MUTED
            p2.font.italic = True

    def _callout(
        self,
        slide,
        title: str,
        lines: list[str],
        left: float = 0.55,
        top: float = 3.85,
        width: float = 6.35,
        height: float = 2.45,
    ) -> None:
        """Light callout for short methodological explanations on a slide."""
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.color.rgb = BORDER
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.16)
        tf.margin_right = Inches(0.14)
        tf.margin_top = Inches(0.1)
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(12)
        p0.font.bold = True
        p0.font.color.rgb = PRIMARY
        p0.space_after = Pt(4)
        for line in lines:
            p = tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(10.5)
            p.font.color.rgb = DARK
            p.space_after = Pt(3)

    def _takeaway(self, slide, text: str) -> None:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(6.5), Inches(12.2), Inches(0.55)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.color.rgb = ACCENT
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.18)
        tf.margin_right = Inches(0.18)
        tf.margin_top = Inches(0.04)
        p = tf.paragraphs[0]
        p.text = f"À retenir — {text}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = PRIMARY

    def _bullets(
        self,
        slide,
        items: list[str],
        left: float = 0.65,
        top: float = BODY_TOP,
        width: float = 6.0,
        height: float = 4.7,
        size: int = 13,
    ) -> None:
        body = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        body.line.fill.background()
        tf = body.text_frame
        tf.word_wrap = True
        for i, line in enumerate(items):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = f"• {line}"
            para.font.size = Pt(size)
            para.font.color.rgb = DARK
            para.space_after = Pt(5)

    def _add_image(self, slide, path: Path, left: float, top: float, width: float) -> bool:
        if path.exists():
            slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
            self.plots_used.append(str(path.relative_to(ROOT)))
            return True
        self.plots_missing.append(str(path.relative_to(ROOT)))
        ph = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(width * 0.6)
        )
        ph.fill.solid()
        ph.fill.fore_color.rgb = LIGHT_BG
        ph.line.color.rgb = BORDER
        ph.text_frame.text = f"[Missing: {path.name}]"
        return False

    def _section_block(
        self,
        slide,
        sections: list[tuple[str, str]],
        left: float,
        top: float,
        width: float,
        title_size: int = 11,
        body_size: int = 11.5,
        spacing: int = 10,
        label_color: RGBColor | None = None,
        height: float = 4.9,
    ) -> None:
        """Render a stack of (label, paragraph) sections without bullets."""
        lc = label_color if label_color is not None else ACCENT
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for label, body in sections:
            p_label = tf.paragraphs[0] if first else tf.add_paragraph()
            p_label.text = label.upper()
            p_label.font.size = Pt(title_size)
            p_label.font.bold = True
            p_label.font.color.rgb = lc
            p_label.space_before = Pt(0 if first else spacing)
            first = False
            p_body = tf.add_paragraph()
            p_body.text = body
            p_body.font.size = Pt(body_size)
            p_body.font.color.rgb = DARK
            p_body.space_before = Pt(2)

    def _stat_chips(
        self,
        slide,
        stats: list[tuple[str, str]],
        left: float = 0.55,
        top: float = 1.55,
        width: float = 12.2,
        height: float = 1.05,
    ) -> None:
        """Horizontal row of large hero numbers + labels (no card backgrounds)."""
        n = len(stats)
        col_w = width / n
        for i, (value, label) in enumerate(stats):
            x = left + i * col_w
            # Value (big number)
            tb_val = slide.shapes.add_textbox(Inches(x), Inches(top), Inches(col_w), Inches(0.50))
            tb_val.line.fill.background()
            tf = tb_val.text_frame
            tf.word_wrap = True
            p_v = tf.paragraphs[0]
            p_v.text = value
            p_v.font.size = Pt(28)
            p_v.font.bold = True
            p_v.font.color.rgb = PRIMARY
            p_v.alignment = PP_ALIGN.CENTER
            # accent underline bar — right below the value box
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + col_w / 2 - 0.35),
                Inches(top + 0.52),
                Inches(0.7),
                Inches(0.04),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = ACCENT
            bar.line.fill.background()
            # Label — separate text box below the bar so it never overlaps
            tb_lbl = slide.shapes.add_textbox(Inches(x), Inches(top + 0.60), Inches(col_w), Inches(0.35))
            tb_lbl.line.fill.background()
            tf_l = tb_lbl.text_frame
            tf_l.word_wrap = True
            p_l = tf_l.paragraphs[0]
            p_l.text = label
            p_l.font.size = Pt(11)
            p_l.font.color.rgb = MUTED
            p_l.alignment = PP_ALIGN.CENTER

    def _class_card(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        title: str,
        threshold: str,
        share: str,
        accent_color: RGBColor,
    ) -> None:
        """Compact card describing one of the 3 outcome classes."""
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.06)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_color
        bar.line.fill.background()
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.16)
        tf.margin_right = Inches(0.16)
        tf.margin_top = Inches(0.16)
        p0 = tf.paragraphs[0]
        p0.text = title.upper()
        p0.font.size = Pt(12)
        p0.font.bold = True
        p0.font.color.rgb = accent_color
        p1 = tf.add_paragraph()
        p1.text = threshold
        p1.font.size = Pt(13)
        p1.font.color.rgb = DARK
        p1.space_before = Pt(4)
        p2 = tf.add_paragraph()
        p2.text = share
        p2.font.size = Pt(10)
        p2.font.color.rgb = MUTED
        p2.space_before = Pt(4)

    def _col_divider(self, slide, top: float = 1.45, height: float = 4.85) -> None:
        """Thin vertical separator between left and right columns."""
        div = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.60),
            Inches(top),
            Inches(0.02),
            Inches(height),
        )
        div.fill.solid()
        div.fill.fore_color.rgb = BORDER
        div.line.fill.background()

    def _metrics_sidebar(
        self,
        slide,
        metrics: list[tuple[str, str]],
        left: float = 7.05,
        top: float = 1.65,
        width: float = 5.75,
    ) -> None:
        """Right-column KPIs as plain text (no white cards that cover neighbouring content)."""
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.5))
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for label, value in metrics:
            pv = tf.paragraphs[0] if first else tf.add_paragraph()
            pv.text = value
            pv.font.size = Pt(30)
            pv.font.bold = True
            pv.font.color.rgb = PRIMARY
            pv.space_before = Pt(0 if first else 16)
            first = False
            pl = tf.add_paragraph()
            pl.text = label
            pl.font.size = Pt(11)
            pl.font.color.rgb = MUTED
            pl.space_after = Pt(4)

    # ----- slides -----

    def slide_01_title(self) -> None:
        slide = self._new_slide()
        # left band
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.55), SLIDE_H)
        band.fill.solid()
        band.fill.fore_color.rgb = PRIMARY
        band.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(11.5), Inches(4.0))
        tb.line.fill.background()
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Prédire le succès d'un film avant sa sortie"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = DARK

        p2 = tf.add_paragraph()
        p2.text = "Classification flop / moyen / hit à partir des données TMDB 5000"
        p2.font.size = Pt(18)
        p2.font.color.rgb = MUTED
        p2.space_before = Pt(14)

        p3 = tf.add_paragraph()
        p3.text = "Projet machine learning · 7 notebooks · application Streamlit"
        p3.font.size = Pt(13)
        p3.font.color.rgb = ACCENT
        p3.space_before = Pt(16)

        # Author block
        p4 = tf.add_paragraph()
        p4.text = "Paul-César Bensemoun"
        p4.font.size = Pt(20)
        p4.font.bold = True
        p4.font.color.rgb = PRIMARY
        p4.space_before = Pt(36)

        p5 = tf.add_paragraph()
        p5.text = "Albert School — ML Proof of Concept"
        p5.font.size = Pt(14)
        p5.font.color.rgb = MUTED
        p5.space_before = Pt(4)

        # bottom band
        bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.0), SLIDE_W, Inches(0.5))
        bot.fill.solid()
        bot.fill.fore_color.rgb = LIGHT_BG
        bot.line.fill.background()
        cap = slide.shapes.add_textbox(Inches(1.0), Inches(7.1), Inches(11.0), Inches(0.3))
        cap.line.fill.background()
        cap.text_frame.text = "Présentation de soutenance · dataset TMDB 5000 movies + credits"
        cap.text_frame.paragraphs[0].font.size = Pt(10)
        cap.text_frame.paragraphs[0].font.color.rgb = MUTED

    def slide_02_data(self) -> None:
        """Données utilisées : TMDB 5000."""
        slide = self._new_slide()
        self._accent_bar(slide)
        self._title(
            slide,
            "Le dataset TMDB 5000",
            "Source publique de référence pour le ML cinéma",
        )

        self._stat_chips(
            slide,
            [
                ("4 803", "films au départ"),
                ("2", "fichiers (movies + credits)"),
                ("1916 — 2017", "période couverte"),
                ("~20", "variables exploitables"),
            ],
            top=1.65,
        )

        self._col_divider(slide, top=3.30, height=3.05)

        self._section_block(
            slide,
            [
                (
                    "Ce qu'on a",
                    "Movies — budget, durée, genre, langue, date de sortie, sociétés et pays de "
                    "production, revenu. Credits — pour chaque film, la liste complète du cast "
                    "et de l'équipe technique au format JSON.",
                ),
            ],
            left=0.55,
            top=3.35,
            width=6.0,
            height=2.90,
        )

        self._section_block(
            slide,
            [
                (
                    "Ce qu'on n'a pas",
                    "Aucun signal marché pré-sortie : budget marketing, critiques, buzz réseaux "
                    "sociaux, calendrier concurrentiel. Conséquence : on modélise uniquement avec "
                    "des variables disponibles avant le tournage (politique anti-leakage stricte).",
                ),
            ],
            left=6.85,
            top=3.35,
            width=6.0,
            height=2.90,
        )

        self._takeaway(
            slide,
            "TMDB est suffisant pour une preuve de concept pré-sortie, mais ne contient pas les "
            "signaux marché — c'est une limite assumée du projet.",
        )

    def slide_03_audit_target(self) -> None:
        """Notebook 01 — audit + définition de la cible."""
        slide = self._new_slide()
        self._accent_bar(slide)
        self._title(
            slide,
            "Nettoyage des données et définition du succès",
            "Du dataset brut à un dataset de modélisation propre",
        )

        self._section_block(
            slide,
            [
                (
                    "Nettoyage",
                    "Parsing des colonnes JSON (genres, sociétés, keywords, cast, crew) puis "
                    "suppression des films sans budget ni revenu exploitables — 1 574 films "
                    "retirés sur 4 803. Résultat : 3 229 films prêts à modéliser.",
                ),
                (
                    "Cible",
                    "On calcule le ROI (revenu ÷ budget) uniquement pour étiqueter les films, "
                    "jamais comme variable d'entrée. Trois classes définies par des seuils "
                    "explicites — voir cartes ci-dessous.",
                ),
                (
                    "Anti-leakage",
                    "Revenu, ROI, popularité et votes sont strictement interdits dans les "
                    "features : on ne modélise qu'avec des signaux disponibles avant la sortie.",
                ),
            ],
            left=0.55,
            top=1.55,
            width=6.2,
            spacing=8,
        )

        self._col_divider(slide, top=1.45, height=3.60)
        self._add_image(slide, self._plot("02_success_class_counts.png", "business_eda"), 7.05, 1.55, 5.7)

        # Three class cards across the bottom
        card_w = 3.95
        card_h = 1.15
        card_top = 5.20
        self._class_card(slide, 0.55, card_top, card_w, card_h,
                         "Flop", "ROI < 1×", "24 % des films", RGBColor(0xB4, 0x1F, 0x1F))
        self._class_card(slide, 4.69, card_top, card_w, card_h,
                         "Moyen", "1 ≤ ROI < 2×", "20 % des films", RGBColor(0xB4, 0x6D, 0x1F))
        self._class_card(slide, 8.83, card_top, card_w, card_h,
                         "Hit", "ROI ≥ 2×", "56 % des films", SUCCESS)

        self._takeaway(
            slide,
            "Trois classes définies par des seuils ROI explicites — chaque film est étiqueté de "
            "façon objective et reproductible.",
        )

    def slide_04_eda(self) -> None:
        """Notebook 02 — EDA."""
        slide = self._new_slide()
        self._accent_bar(slide)
        self._title(
            slide,
            "Premières observations sur les données",
            "Quelles variables semblent contenir du signal ?",
        )

        self._section_block(
            slide,
            [
                (
                    "Budget",
                    "Le budget seul ne suffit pas : on trouve des hits à 5 M$ et des flops à "
                    "200 M$. Sa distribution est extrêmement asymétrique → échelle log "
                    "indispensable pour modéliser.",
                ),
                (
                    "Genre",
                    "Profils ROI très contrastés : horreur et animation surperforment, drame et "
                    "documentaire sont plus risqués.",
                ),
                (
                    "Timing",
                    "Saisonnalité visible : les sorties estivales et de fin d'année concentrent "
                    "les meilleurs ROI.",
                ),
                (
                    "Conclusion",
                    "Aucune variable seule n'explique le succès — il faut combiner et créer des "
                    "signaux dérivés (feature engineering).",
                ),
            ],
            left=0.55,
            top=1.55,
            width=6.2,
            spacing=9,
        )

        self._col_divider(slide)
        self._add_image(slide, self._plot("04b_logroi_by_main_genre.png", "business_eda"), 7.05, 1.6, 5.7)

        self._takeaway(
            slide,
            "Trois axes porteurs : budget, genre et timing. Cela oriente le feature engineering "
            "qui suit.",
        )

    def slide_05_baseline(self) -> None:
        """Notebook 03 — modeling baseline."""
        slide = self._new_slide()
        self._accent_bar(slide)
        self._title(
            slide,
            "Modèle de référence (baseline)",
            "Un premier modèle pour fixer un point de comparaison",
        )

        self._section_block(
            slide,
            [
                (
                    "Setup",
                    "10 variables de base (budget, durée, genre, langue, mois de sortie). "
                    "Split train / test 80 / 20 stratifié sur les 3 classes.",
                ),
                (
                    "Modèles comparés",
                    "Régression logistique, forêt aléatoire, gradient boosting. Champion retenu "
                    "= meilleur macro-F1 sur le test → régression logistique (macro-F1 ≈ 0.37, "
                    "accuracy ≈ 39 %).",
                ),
            ],
            left=0.55,
            top=1.55,
            width=6.2,
            spacing=10,
        )

        self._callout(
            slide,
            "Pourquoi macro-F1 ?",
            [
                "Classes déséquilibrées (~56 % hit, ~24 % flop, ~20 % moyen) : l'accuracy "
                "récompense « tout prédire hit ».",
                "Macro-F1 = moyenne des F1 des 3 classes → chaque classe compte autant.",
                "C'est aligné avec l'objectif métier : détecter flop, moyen ET hit.",
            ],
            left=0.55,
            top=4.10,
            width=6.20,
            height=2.20,
        )

        self._col_divider(slide)
        self._add_image(slide, self._plot("01_model_comparison_metrics.png", "modeling"), 7.05, 1.65, 5.7)

        self._takeaway(
            slide,
            "On compare les modèles sur macro-F1 — métrique honnête face au déséquilibre. "
            "Le baseline à 0.37 sert de référence pour toute la suite du projet.",
        )

    def slide_06_feature_engineering(self) -> None:
        """Notebook 04 — feature engineering."""
        slide = self._new_slide()
        self._accent_bar(slide)
        self._title(
            slide,
            "Feature engineering : créer de nouvelles variables",
            "Enrichir le modèle sans ajouter de nouvelles données",
        )

        # Left column — short paragraphs (no bullets)
        self._section_block(
            slide,
            [
                (
                    "Transformations clés",
                    "Budget en échelle logarithmique pour corriger l'asymétrie extrême "
                    "(1 K$ → 380 M$). Durée bucketisée (court / moyen / long). Décennie de "
                    "sortie comme proxy d'époque.",
                ),
                (
                    "Signaux composites",
                    "Complexité du genre (focused / mixed / hybrid) et échelle de production "
                    "(indie / mid / large) calculées sur les tertiles du train uniquement — "
                    "anti-fuite vérifié après chaque transformation.",
                ),
                (
                    "Gain",
                    "Régression logistique enrichie : macro-F1 ≈ 0.39 (+0.02 vs baseline).",
                ),
            ],
            left=0.55,
            top=1.55,
            width=6.05,
            spacing=8,
        )

        self._col_divider(slide)

        # Right column — feature list aligned with production champion (25 features)
        right = slide.shapes.add_textbox(Inches(6.85), Inches(1.55), Inches(5.95), Inches(4.7))
        right.line.fill.background()
        tf = right.text_frame
        tf.word_wrap = True
        head = tf.paragraphs[0]
        head.text = "LES 25 VARIABLES DU MODÈLE DE PRODUCTION"
        head.font.size = Pt(11)
        head.font.bold = True
        head.font.color.rgb = PRIMARY

        groups = [
            ("Budget & durée (3)",
             "budget_log · runtime · runtime_bucket"),
            ("Genre & langue (5)",
             "main_genre · genre_count · genre_complexity · original_language · spoken_language_count"),
            ("Production & timing (7)",
             "production_company_count · production_country_count · international_production · multilingual_movie · production_scale · decade · release_month"),
            ("Casting (4)",
             "cast_size · top_billed_cast_count · known_actor_count · ensemble_cast_flag"),
            ("Équipe & réalisateur (4)",
             "crew_size · writer_count · director_movie_count · top_director_flag"),
            ("Signaux composites (2)",
             "talent_score · possible_franchise_flag"),
        ]
        for label, items in groups:
            p_lab = tf.add_paragraph()
            p_lab.text = label
            p_lab.font.size = Pt(10)
            p_lab.font.bold = True
            p_lab.font.color.rgb = ACCENT
            p_lab.space_before = Pt(4)
            p_items = tf.add_paragraph()
            p_items.text = items
            p_items.font.size = Pt(9)
            p_items.font.color.rgb = MUTED

        self._takeaway(
            slide,
            "Le feature engineering apporte +0.02 de macro-F1. Les variables credits "
            "ajoutées dans le notebook 05 feront le plus gros bond.",
        )

    def slide_07_credits(self) -> None:
        """Notebook 05 — credits enrichment."""
        slide = self._new_slide()
        self._accent_bar(slide)
        self._title(
            slide,
            "Enrichissement avec les credits : le saut de qualité",
            "On intègre le cast, l'équipe technique et la dimension franchise",
        )

        self._section_block(
            slide,
            [
                (
                    "Signaux ajoutés",
                    "Cast (taille, acteurs connus, top billed, ensemble flag), équipe (crew "
                    "size, scénaristes), réalisateur (nombre de films + flag top réalisateur), "
                    "franchise (heuristique titre + keywords).",
                ),
                (
                    "Score talent composite",
                    "Combinaison pondérée des signaux acteurs / réalisateur pour résumer "
                    "le « packaging créatif » en une seule variable.",
                ),
                (
                    "Top features (permutation importance)",
                    "Taille du cast, décennie, expérience du réalisateur, nombre de sociétés "
                    "de production et budget (log) sont les leviers les plus influents — "
                    "des signaux de packaging et d'industrie, pas de simples métadonnées.",
                ),
                (
                    "Gain",
                    "Régime credits — macro-F1 ≈ 0.44 (+0.05 vs FE, +0.07 vs baseline) · "
                    "accuracy ≈ 47 %.",
                ),
            ],
            left=0.55,
            top=1.55,
            width=6.05,
            spacing=8,
        )

        self._col_divider(slide)
        self._add_image(
            slide,
            self._plot("30_clean_champion_feature_importance.png", "modeling"),
            6.85,
            1.55,
            6.0,
        )

        self._takeaway(
            slide,
            "Le plus gros gain du projet vient des signaux d'industrie et de packaging — "
            "cast, réalisateur, sociétés de production, budget et franchise.",
        )

    def slide_08_error_analysis(self) -> None:
        """Notebook 06 — error analysis."""
        slide = self._new_slide()
        self._accent_bar(slide)
        self._title(
            slide,
            "Analyse des erreurs : qu'a apporté l'enrichissement credits ?",
            "On compare les prédictions baseline vs credits sur le même test set",
        )

        n_w2c = int(self.kpis.get("n_wrong_to_correct", 130))
        n_c2w = int(self.kpis.get("n_correct_to_wrong", 79))
        n_test = int(self.kpis.get("n_test", 646))

        self._section_block(
            slide,
            [
                (
                    "Méthode",
                    f"On compare film par film les prédictions du baseline et du modèle credits "
                    f"sur le même test set ({n_test} films) et on classe chaque film dans "
                    f"l'une des quatre situations ci-dessous.",
                ),
                (
                    "Sauvés par les credits",
                    f"{n_w2c} films que le baseline ratait sont désormais correctement classés "
                    f"grâce aux signaux talent / packaging.",
                ),
                (
                    "Perdus par les credits",
                    f"{n_c2w} films sont devenus faux après enrichissement — quelques signaux "
                    f"talent peuvent pousser le modèle dans la mauvaise direction.",
                ),
                (
                    "Gain net",
                    f"+{n_w2c - n_c2w} prédictions correctes. Le solde est positif : "
                    f"l'enrichissement crée plus de valeur qu'il n'en détruit.",
                ),
                (
                    "Limite structurelle",
                    "La classe « moyen » (films qui rendent entre 1× et 2× leur budget) reste "
                    "la plus difficile pour les deux modèles — c'est une fenêtre ROI étroite, "
                    "pas un défaut du modèle.",
                ),
            ],
            left=0.55,
            top=1.55,
            width=6.05,
            spacing=7,
        )

        self._col_divider(slide)
        self._add_image(slide, self._plot("12_error_transition_counts.png", "modeling"), 6.85, 1.55, 6.0)

        self._takeaway(
            slide,
            f"L'enrichissement credits apporte un gain net mesurable (+{n_w2c - n_c2w} films "
            f"mieux classés). La classe moyen reste le point dur — limite de la fenêtre ROI.",
        )

    def slide_09_advanced_modeling(self) -> None:
        """Notebook 07 — advanced modeling, expliqué simplement."""
        slide = self._new_slide()
        self._accent_bar(slide)
        self._title(
            slide,
            "Validation et optimisation du modèle final",
            "Quatre techniques pour solidifier le champion",
        )

        self._section_block(
            slide,
            [
                (
                    "Validation croisée (5 plis)",
                    "Macro-F1 = 0.441 ± 0.011 → performance très stable, le score test "
                    "n'est pas dû à la chance.",
                ),
                (
                    "Recherche d'hyperparamètres",
                    "5 familles de modèles testées (LR, RF, GB, XGBoost, LR + overview). "
                    "La régression logistique tunée reste le meilleur compromis "
                    "performance / interprétabilité.",
                ),
                (
                    "Calibration des probabilités",
                    "Calibration isotonique : +10 pp d'accuracy mais −0.08 macro-F1 "
                    "(la classe moyen s'effondre). On garde la version non calibrée "
                    "pour rester utile sur les 3 classes.",
                ),
                (
                    "Approche en cascade (ordinale)",
                    "Modèle qui prédit d'abord flop vs non-flop, puis hit vs moyen. "
                    "Macro-F1 ≈ 0.34 → sous-performe, on conserve l'approche directe.",
                ),
            ],
            left=0.55,
            top=1.55,
            width=6.30,
            spacing=9,
        )

        self._col_divider(slide)
        self._metrics_sidebar(
            slide,
            [
                ("Macro-F1 final (test)", self._pct("champion_macro_f1", 0.4445)),
                (
                    "Macro-F1 validation croisée",
                    f"{float(self.kpis.get('champion_cv_macro_f1_mean', 0.441)):.3f} "
                    f"± {float(self.kpis.get('champion_cv_macro_f1_std', 0.011)):.3f}",
                ),
                ("Accuracy", self._pct("champion_accuracy", 0.478)),
                ("Modèles comparés", "5 familles × 25 réglages"),
            ],
        )

        self._takeaway(
            slide,
            "Performance stable en cross-validation, trade-offs documentés. "
            "On valide le champion et on le garde pour l'application.",
        )

    def slide_10_app_and_conclusion(self) -> None:
        """Conclusion + application Streamlit."""
        slide = self._new_slide()
        self._accent_bar(slide)
        self._title(
            slide,
            "Bilan du projet & application Streamlit",
            "Ce qu'on a produit, les limites, et les suites possibles",
        )

        # Progression banner
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.55), Inches(12.20), Inches(0.60)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = ACCENT_SOFT
        banner.line.color.rgb = ACCENT
        banner.line.width = Pt(0.75)
        tb = slide.shapes.add_textbox(Inches(0.70), Inches(1.60), Inches(12.0), Inches(0.50))
        tb.line.fill.background()
        tf = tb.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = (
            "Progression macro-F1 : 0.373 (baseline)  →  0.390 (feature engineering)  →  "
            "0.438 (credits)  →  0.4445 (tuning + refit sur 25 features nettoyées)"
        )
        p0.font.size = Pt(12)
        p0.font.bold = True
        p0.font.color.rgb = PRIMARY

        self._col_divider(slide, top=2.30, height=4.00)

        self._section_block(
            slide,
            [
                (
                    "Ce qui fonctionne",
                    "Pipeline complet et reproductible (7 notebooks). Discipline anti-leakage "
                    "rigoureuse. Champion : LR tunée sur 25 features nettoyées (multicolinéarité "
                    "résolue). Application Streamlit opérationnelle : prédiction, comparaison "
                    "de scénarios, comparables historiques, export PDF.",
                ),
            ],
            left=0.55,
            top=2.28,
            width=5.90,
            label_color=SUCCESS,
            spacing=0,
            height=4.10,
        )

        self._section_block(
            slide,
            [
                (
                    "Limites",
                    "Pas de données marketing ni de buzz pré-sortie dans TMDB. "
                    "La classe moyen (ROI entre 1× et 2×) reste difficile à prédire — "
                    "fenêtre ROI étroite (F1 ≈ 0.30).",
                ),
                (
                    "Suites possibles",
                    "Intégrer des critiques (Rotten Tomatoes, Metacritic) et données "
                    "Box Office Mojo. Valider en temporel (entraîner sur le passé, "
                    "tester sur le futur).",
                ),
            ],
            left=6.65,
            top=2.28,
            width=6.10,
            spacing=10,
            height=4.10,
        )

        self._takeaway(
            slide,
            "Pipeline ML rigoureux + application crédible pour un usage métier. "
            "La performance est limitée par les données disponibles, pas par la méthodologie.",
        )

    def slide_09_limitations(self) -> None:
        # Kept (deprecated, no longer used in build sequence)
        slide = self._new_slide()
        self._accent_bar(slide)
        self._title(
            slide,
            "Limites & axes d'amélioration",
            "Ce qui pourrait expliquer le plafond actuel et comment progresser",
        )

        # Left column: limits
        col1 = slide.shapes.add_textbox(Inches(0.55), Inches(1.5), Inches(6.0), Inches(4.7))
        tf = col1.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = "Limites identifiées"
        p0.font.size = Pt(15)
        p0.font.bold = True
        p0.font.color.rgb = PRIMARY
        p0.space_after = Pt(8)
        limits = [
            "Pas de variables marketing (budget pub, sorties concurrentes)",
            "Pas de signal audience pré-sortie (buzz social, trailers)",
            "Budget / revenue non ajustés à l'inflation",
            "Biais TMDB : sur-représentation des films grand public anglo-saxons",
            "Classe average reste difficile (F1 0.305 après tuning, 0.016 après calibration)",
            "Random split, pas de validation temporelle (suggérée comme follow-up)",
            "Embeddings sentence-transformer non testés (LSA-50 a sous-performé)",
        ]
        for item in limits:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(12)
            p.font.color.rgb = DARK
            p.space_after = Pt(4)

        # Right column: improvements
        col2 = slide.shapes.add_textbox(Inches(6.85), Inches(1.5), Inches(6.0), Inches(4.7))
        tf = col2.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = "Pistes d'amélioration"
        p0.font.size = Pt(15)
        p0.font.bold = True
        p0.font.color.rgb = PRIMARY
        p0.space_after = Pt(8)
        impros = [
            "✓ Cross-validation 5-fold + RandomizedSearchCV (notebook 07)",
            "✓ Calibration isotonique implémentée (notebook 07)",
            "✓ Ordinal cascade testée vs multinomial (notebook 07)",
            "✓ Overview LSA testée (TF-IDF + TruncatedSVD)",
            "→ Sentence-transformer embeddings (PyTorch, follow-up)",
            "→ Validation temporelle (train < 2014 / test ≥ 2014)",
            "→ Données marché externes (Box Office Mojo, IMDb, Google Trends)",
            "→ Mord (ordinal regression) ou cascade plus profonde",
        ]
        for item in impros:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(12)
            p.font.color.rgb = DARK
            p.space_after = Pt(4)

        self._takeaway(
            slide,
            "Le plafond actuel reflète surtout l'absence de signaux marché — ajouter des features "
            "ML ne suffira pas, il faut élargir la donnée.",
        )

    def slide_10_advanced_modeling(self) -> None:
        """Dedicated slide for the notebook 07 advanced modeling work."""
        slide = self._new_slide()
        self._accent_bar(slide)
        self._title(
            slide,
            "Cross-validation, tuning, calibration & ordinal cascade",
            "Closing the methodology gap from notebooks 03–06",
        )

        self._metrics_sidebar(
            slide,
            [
                ("Macro-F1 lift", "+0.012"),
                ("CV mean ± std", "0.441 ± 0.011"),
                ("Models compared", "5 (LR/RF/GB/XGB/LSA)"),
                ("Tuning iter × folds", "25 × 5"),
            ],
            left=0.55,
            top=1.65,
            width=12.0,
        )

        # Left column: techniques
        col1 = slide.shapes.add_textbox(Inches(0.55), Inches(3.15), Inches(6.0), Inches(3.7))
        tf = col1.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = "Techniques implémentées"
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = PRIMARY
        p0.space_after = Pt(6)
        for line in [
            "StratifiedKFold(5) + cross_val_score → mean ± std macro-F1",
            "RandomizedSearchCV (25 iter) × 5 model families",
            "CalibratedClassifierCV (isotonic, cv=5) sur le champion",
            "Ordinal cascade : flop vs not-flop → hit vs average",
            "Overview LSA : TfidfVectorizer + TruncatedSVD(50)",
        ]:
            p = tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK
            p.space_after = Pt(3)

        # Right column: findings
        col2 = slide.shapes.add_textbox(Inches(6.85), Inches(3.15), Inches(6.0), Inches(3.7))
        tf = col2.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = "Résultats clés"
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = PRIMARY
        p0.space_after = Pt(6)
        for line in [
            "LR tuned : best macro-F1 (0.450) → champion production",
            "GradientBoosting tuned : best accuracy (0.556) + best log-loss",
            "Calibration isotonique : +9 pp accuracy MAIS -0.08 macro-F1",
            "Ordinal cascade sous-performe (-0.10 macro-F1)",
            "LSA-50 overfit : CV +0.004, test -0.024",
        ]:
            p = tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(11)
            p.font.color.rgb = DARK
            p.space_after = Pt(3)

        self._takeaway(
            slide,
            "La rigueur méthodologique (CV, tuning, calibration) est désormais en place. "
            "Le trade-off macro-F1 / accuracy / log-loss est documenté et choisi en conscience.",
        )

    def slide_11_conclusion(self) -> None:
        slide = self._new_slide()
        self._accent_bar(slide)
        self._title(
            slide,
            "Conclusion technique",
            "Bilan du projet et message principal",
        )

        body = slide.shapes.add_textbox(Inches(0.65), Inches(1.6), Inches(12.0), Inches(4.6))
        tf = body.text_frame
        tf.word_wrap = True

        sections = [
            (
                "Ce qui a fonctionné",
                [
                    "Pipeline ML complet, reproductible, leakage-safe (notebooks 01 → 07)",
                    "Progression macro-F1 documentée : 0.373 → 0.390 → 0.438 → 0.450 (tuning)",
                    "Enrichissement crédits valide l'hypothèse : signaux talent / packaging comptent",
                    "Analyse d'erreurs + cross-validation + tuning + calibration — méthodo complète",
                    "Production champion sélectionné par macro-F1 test parmi 5 familles tunées",
                ],
                PRIMARY,
            ),
            (
                "Ce qui n'a pas fonctionné (et c'est intéressant)",
                [
                    "Classe average reste difficile (F1 0.305) — frontière ROI ambiguë",
                    "Overview LSA n'a pas généralisé (overfit) — sentence-transformer à tester",
                    "Calibration isotonique : trade-off accuracy ↑ vs macro-F1 ↓ documenté",
                ],
                MUTED,
            ),
            (
                "Message principal",
                [
                    "Le ML pré-sortie ne prédit pas parfaitement le succès, mais structure la "
                    "décision avec un signal mesurable, explicable, et désormais correctement validé "
                    "par cross-validation et tuning systématique.",
                ],
                ACCENT,
            ),
        ]

        first = True
        for header, lines, color in sections:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.text = header
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = color
            p.space_before = Pt(0 if first else 14)
            first = False
            for line in lines:
                pl = tf.add_paragraph()
                pl.text = f"• {line}"
                pl.font.size = Pt(12)
                pl.font.color.rgb = DARK
                pl.space_after = Pt(3)

        self._takeaway(
            slide,
            "La valeur du projet n'est pas la prédiction parfaite, mais la rigueur méthodologique "
            "et l'analyse business du résultat.",
        )

    def build(self) -> None:
        self.slide_01_title()
        self.slide_02_data()
        self.slide_03_audit_target()
        self.slide_04_eda()
        self.slide_05_baseline()
        self.slide_06_feature_engineering()
        self.slide_07_credits()
        self.slide_08_error_analysis()
        self.slide_09_advanced_modeling()
        self.slide_10_app_and_conclusion()

    def save(self) -> None:
        PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(OUTPUT_PPTX))


def main() -> None:
    deck = TechDeck()
    deck.build()
    deck.save()

    print("===== TECHNICAL DECK BUILD =====")
    print(f"Output: {OUTPUT_PPTX}")
    print(f"Slides: {len(deck.prs.slides)}")
    print(f"Plots embedded ({len(deck.plots_used)}):")
    for p in deck.plots_used:
        print(f"  + {p}")
    if deck.plots_missing:
        print(f"Plots missing ({len(deck.plots_missing)}):")
        for p in deck.plots_missing:
            print(f"  - {p}")


if __name__ == "__main__":
    main()
